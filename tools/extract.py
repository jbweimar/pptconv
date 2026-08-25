#!/usr/bin/env python
"""Extract a deck's full content into a manifest (work/content.json) + assets (work/assets/).

The manifest is the neutral middle format of the conversion: everything the rebuild
step needs, decoupled from the old template's formatting. Text keeps structure
(paragraph levels) and semantic emphasis only (bold/italic runs); all other
formatting is deliberately dropped so the NEW template's styles rule.

Usage: extract.py deck.pptx [outdir=work]
"""
from __future__ import annotations
import hashlib
import json
import sys
from pathlib import Path

from pptx import Presentation


def runs_of(paragraph):
    runs = []
    for r in paragraph.runs:
        runs.append({
            "text": r.text,
            **({"bold": True} if r.font.bold else {}),
            **({"italic": True} if r.font.italic else {}),
            **({"hyperlink": r.hyperlink.address} if r.hyperlink and r.hyperlink.address else {}),
        })
    return runs


def text_frame_of(tf):
    paras = []
    for p in tf.paragraphs:
        paras.append({"level": p.level, "runs": runs_of(p), "text": "".join(r.text for r in p.runs)})
    return paras


def save_image(shape, assets: Path):
    img = shape.image
    digest = hashlib.sha1(img.blob).hexdigest()[:12]
    ext = img.ext
    fname = f"{digest}.{ext}"
    path = assets / fname
    if not path.exists():
        path.write_bytes(img.blob)
    return {"file": f"assets/{fname}", "size_px": list(img.size), "content_type": img.content_type}


def table_of(table):
    return {
        "rows": [[text_frame_of(cell.text_frame) for cell in row.cells] for row in table.rows],
        "n_rows": len(table.rows), "n_cols": len(table.columns),
    }


def chart_of(chart):
    try:
        plots = []
        cats = [str(c) for c in chart.plots[0].categories] if chart.plots else []
        for plot in chart.plots:
            for s in plot.series:
                plots.append({"name": s.name, "values": [v for v in s.values]})
        return {"chart_type": str(chart.chart_type), "categories": cats, "series": plots}
    except Exception as e:
        return {"chart_type": "UNREADABLE", "error": str(e)}


def shape_content(shape, assets: Path):
    d = {"name": shape.name}
    if shape.is_placeholder:
        pf = shape.placeholder_format
        d["placeholder"] = {"idx": pf.idx,
                            "type": str(pf.type).split(".")[-1].split(" ")[0] if pf.type is not None else None}
    st = str(shape.shape_type)
    if "PICTURE" in st:
        d["image"] = save_image(shape, assets)
    elif getattr(shape, "has_table", False):
        d["table"] = table_of(shape.table)
    elif getattr(shape, "has_chart", False):
        d["chart"] = chart_of(shape.chart)
    elif "GROUP" in st:
        d["group"] = [shape_content(s, assets) for s in shape.shapes]
        d["review"] = "grouped shapes: verify manually after rebuild"
    elif shape.has_text_frame and shape.text_frame.text.strip():
        d["paragraphs"] = text_frame_of(shape.text_frame)
    else:
        return None  # decorative/empty shape from the old template: drop
    if shape.left is not None:
        d["pos_emu"] = [shape.left, shape.top, shape.width, shape.height]
    return d


def extract(path: str, outdir: str = "work"):
    prs = Presentation(path)
    out = Path(outdir)
    assets = out / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    doc = {"source": str(Path(path).name),
           "slide_size_emu": [prs.slide_width, prs.slide_height],
           "slides": []}
    for si, slide in enumerate(prs.slides, 1):
        contents = [c for c in (shape_content(sh, assets) for sh in slide.shapes) if c]
        s = {"n": si, "layout": slide.slide_layout.name, "shapes": contents}
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                s["notes"] = notes
        doc["slides"].append(s)
    manifest = out / "content.json"
    manifest.write_text(json.dumps(doc, indent=2, ensure_ascii=False))
    n_img = len(list(assets.glob("*")))
    print(f"extracted {len(doc['slides'])} slides -> {manifest} ({n_img} unique assets)")
    return doc


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    extract(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "work")
