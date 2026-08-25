#!/usr/bin/env python
"""Project-specific conversion: graft the corporate roadmap slides into the new template.

The roadmap slides are vector diagrams (S-curve timelines) drawn in THEME colors,
so transplanting their shape XML into the new deck re-resolves every scheme color
and theme font against the new template — the diagrams restyle themselves.

Per slide: new slide from the new template's "Title Only" layout (new master bg,
title style), title text moved into the new title placeholder, every other shape
deep-copied. Copied placeholder shapes are DE-PLACEHOLDERED (ph element stripped)
so they become plain shapes; geometry missing on the shape is inherited from the
source layout's matching placeholder first.

Usage: graft_roadmap.py <new_template.pptx> <roadmap.pptx> <out.pptx>
"""
from __future__ import annotations
import copy
import sys

from pptx import Presentation
from pptx.oxml.ns import qn

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def strip_existing_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)


def layout_by_name(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(name)


def ph_of(sp):
    nv = sp.find(qn("p:nvSpPr")) if sp.tag == qn("p:sp") else None
    if nv is None:
        return None
    return nv.find(qn("p:nvPr") + "/" + qn("p:ph"))


def ph_key(ph):
    return (ph.get("type") or "body", ph.get("idx") or "0")


def layout_geometry(layout):
    """Map (ph type, idx) -> deepcopy-able xfrm element from the layout."""
    geo = {}
    for sp in layout.shapes:
        if not sp.is_placeholder:
            continue
        el = sp._element
        xfrm = el.find(qn("p:spPr") + "/" + qn("a:xfrm"))
        ph = ph_of(el)
        if ph is not None and xfrm is not None:
            geo[ph_key(ph)] = xfrm
    return geo


def layout_liststyles(layout):
    """Map (ph type, idx) -> the layout placeholder's lstStyle (fonts/sizes/colors).

    De-placeholdered shapes lose this inheritance, so we bake it into the copied
    shape's own txBody. Scheme-color references inside still re-resolve against
    the DESTINATION theme, which is exactly what we want.
    """
    styles = {}
    for sp in layout.shapes:
        if not sp.is_placeholder:
            continue
        el = sp._element
        ph = ph_of(el)
        lst = el.find(qn("p:txBody") + "/" + qn("a:lstStyle"))
        if ph is not None and lst is not None and len(lst) > 0:
            styles[ph_key(ph)] = lst
    return styles


def merge_liststyle(txBody, layout_lst):
    """Give the shape its layout's lstStyle unless it already defines one."""
    own = txBody.find(qn("a:lstStyle"))
    if own is not None and len(own) > 0:
        return
    if own is not None:
        txBody.remove(own)
    new = copy.deepcopy(layout_lst)
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.addnext(new)


def graft(new_template, roadmap, out_path):
    prs = Presentation(new_template)
    strip_existing_slides(prs)
    src = Presentation(roadmap)
    target_layout = layout_by_name(prs, "Title Only")

    for s_i, s in enumerate(src.slides, 1):
        slide = prs.slides.add_slide(target_layout)
        src_geo = layout_geometry(s.slide_layout)
        src_styles = layout_liststyles(s.slide_layout)

        # title into the new template's title placeholder (inherits new styling)
        title_text = None
        for ph in s.placeholders:
            if "TITLE" in str(ph.placeholder_format.type) and ph.has_text_frame:
                title_text = ph.text_frame.text
        if title_text and slide.shapes.title is not None:
            slide.shapes.title.text_frame.text = title_text

        spTree = slide.shapes._spTree
        for shape in s.shapes:
            el = shape._element
            ph = ph_of(el)
            if ph is not None:
                t = ph.get("type") or "body"
                if t in ("title", "ctrTitle"):
                    continue  # moved above
                if t in ("sldNum", "ftr", "dt"):
                    continue  # plumbing: the new layout supplies its own
            new_el = copy.deepcopy(el)
            new_ph = ph_of(new_el)
            if new_ph is not None:
                # inherit geometry + text styles from the source layout, then de-placeholder
                spPr = new_el.find(qn("p:spPr"))
                if spPr is not None and spPr.find(qn("a:xfrm")) is None:
                    lay_xfrm = src_geo.get(ph_key(new_ph))
                    if lay_xfrm is not None:
                        spPr.insert(0, copy.deepcopy(lay_xfrm))
                    else:
                        print(f"  slide {s_i}: WARNING no geometry for ph {ph_key(new_ph)}")
                lay_lst = src_styles.get(ph_key(new_ph))
                txBody = new_el.find(qn("p:txBody"))
                if lay_lst is not None and txBody is not None:
                    merge_liststyle(txBody, lay_lst)
                new_ph.getparent().remove(new_ph)
            spTree.append(new_el)

        if s.has_notes_slide:
            notes = s.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
        print(f"slide {s_i}: grafted ({len(s.shapes)} src shapes)")

    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        sys.exit(__doc__)
    graft(*sys.argv[1:])
