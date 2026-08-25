#!/usr/bin/env python
"""Inspect a .pptx: slide size, masters, layouts, placeholder inventory, theme, per-slide contents.

Usage:
  inspect.py deck.pptx            human-readable report
  inspect.py deck.pptx --json     machine-readable (for building layout mappings)
"""
from __future__ import annotations
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_PER_IN = 914400


def emu_in(v):
    return round(Emu(v).inches, 2) if v is not None else None


def ph_info(shape):
    pf = shape.placeholder_format
    return {
        "idx": pf.idx,
        "type": str(pf.type).split(".")[-1].split(" ")[0] if pf.type is not None else None,
        "name": shape.name,
    }


def shape_summary(shape):
    d = {"name": shape.name, "type": str(shape.shape_type).split(" ")[0]}
    if shape.is_placeholder:
        d["placeholder"] = ph_info(shape)
    if shape.has_text_frame:
        text = shape.text_frame.text
        d["chars"] = len(text)
        d["text_preview"] = text[:80].replace("\n", " ⏎ ")
    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
        d["picture"] = True
    if getattr(shape, "has_table", False):
        d["table"] = {"rows": len(shape.table.rows), "cols": len(shape.table.columns)}
    if getattr(shape, "has_chart", False):
        d["chart"] = str(shape.chart.chart_type)
    d["pos_in"] = [emu_in(shape.left), emu_in(shape.top), emu_in(shape.width), emu_in(shape.height)]
    return d


def theme_summary(prs):
    # Read theme fonts/colors from the first master's theme part.
    try:
        theme = prs.slide_masters[0].element.getroottree()
        part = prs.slide_masters[0].part
        for rel in part.rels.values():
            if "theme" in rel.reltype:
                tx = rel.target_part.blob.decode("utf-8", "ignore")
                import re
                fonts = re.findall(r'<a:(?:latin|majorFont|minorFont)[^>]*typeface="([^"]+)"', tx)
                colors = re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"/>', tx)[:12]
                return {"typefaces": sorted(set(f for f in fonts if f and not f.startswith("+"))),
                        "srgb_colors_sample": colors}
    except Exception as e:  # theme reading is best-effort
        return {"error": str(e)}
    return {}


def inspect(path: str):
    prs = Presentation(path)
    out = {
        "file": str(Path(path).name),
        "slide_size_in": [emu_in(prs.slide_width), emu_in(prs.slide_height)],
        "aspect": round(prs.slide_width / prs.slide_height, 3),
        "theme": theme_summary(prs),
        "masters": [],
        "slides": [],
    }
    for mi, master in enumerate(prs.slide_masters):
        m = {"index": mi, "name": master.name, "layouts": []}
        for li, layout in enumerate(master.slide_layouts):
            m["layouts"].append({
                "index": li,
                "name": layout.name,
                "placeholders": [ph_info(ph) for ph in layout.placeholders],
            })
        out["masters"].append(m)
    for si, slide in enumerate(prs.slides, 1):
        s = {
            "n": si,
            "layout": slide.slide_layout.name,
            "shapes": [shape_summary(sh) for sh in slide.shapes],
            "has_notes": bool(slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()),
        }
        out["slides"].append(s)
    return out


def human(out):
    w, h = out["slide_size_in"]
    print(f"== {out['file']} — {w}x{h}\" (aspect {out['aspect']}) ==")
    th = out.get("theme") or {}
    if th.get("typefaces"):
        print(f"theme typefaces: {', '.join(th['typefaces'])}")
    for m in out["masters"]:
        print(f"\nmaster [{m['index']}] {m['name']}")
        for l in m["layouts"]:
            phs = ", ".join(f"{p['idx']}:{p['type']}" for p in l["placeholders"])
            print(f"  layout [{l['index']:2}] {l['name']:<28} ph: {phs}")
    print(f"\n{len(out['slides'])} slides:")
    for s in out["slides"]:
        kinds = []
        for sh in s["shapes"]:
            if "picture" in sh:
                kinds.append("img")
            if "table" in sh:
                kinds.append("tbl")
            if "chart" in sh:
                kinds.append("cht")
        extra = ("  [" + ",".join(kinds) + "]") if kinds else ""
        notes = "  +notes" if s["has_notes"] else ""
        title = next((sh.get("text_preview", "") for sh in s["shapes"]
                      if sh.get("placeholder", {}).get("type") in ("TITLE", "CENTER_TITLE")), "")
        print(f"  {s['n']:3}. layout={s['layout']:<28} {title[:50]}{extra}{notes}")


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        sys.exit(__doc__)
    result = inspect(args[0])
    if "--json" in sys.argv:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    else:
        human(result)
