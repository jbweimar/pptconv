#!/usr/bin/env python
"""Rebuild a deck onto the NEW template from the extracted manifest.

Strategy: REBUILD (not retrofit) — open the new template, add fresh slides from its
layouts, pour content into placeholders. The new template's styles rule; only
structure and semantic emphasis survive from the old deck (see docs/REFERENCE.md).

Usage: apply.py <new_template.pptx> <work/content.json> <mapping.json> <out.pptx>

mapping.json:
{
  "default_layout": "Title and Content",
  "layouts": { "<old layout name>": "<new layout name>", ... },
  "drop_slides": [],            // 1-based slide numbers to skip
  "notes": "free-form comments"
}
"""
from __future__ import annotations
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REVIEW = []


def flag(slide_n, msg):
    REVIEW.append(f"slide {slide_n}: {msg}")


def layout_by_name(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def fill_text_frame(tf, paragraphs):
    tf.clear()
    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = min(p.get("level", 0), 8)
        runs = p.get("runs") or ([{"text": p.get("text", "")}] if p.get("text") else [])
        for r in runs:
            run = para.add_run()
            run.text = r.get("text", "")
            if r.get("bold"):
                run.font.bold = True
            if r.get("italic"):
                run.font.italic = True
            if r.get("hyperlink"):
                run.hyperlink.address = r["hyperlink"]


def body_placeholders(slide):
    """Content-capable placeholders in layout order, excluding title/footer plumbing."""
    skip = {"TITLE", "CENTER_TITLE", "SLIDE_NUMBER", "FOOTER", "DATE"}
    out = []
    for ph in slide.placeholders:
        t = str(ph.placeholder_format.type).split(".")[-1].split(" ")[0] if ph.placeholder_format.type is not None else None
        if t not in skip:
            out.append(ph)
    return out


def title_placeholder(slide):
    for ph in slide.placeholders:
        t = str(ph.placeholder_format.type)
        if "TITLE" in t:
            return ph
    return None


def content_area(prs):
    """Fallback geometry when a shape has no placeholder home: centered, 70% of slide."""
    w, h = prs.slide_width, prs.slide_height
    return Emu(int(w * .15)), Emu(int(h * .2)), Emu(int(w * .7)), Emu(int(h * .65))


def place_image(slide, prs, img_entry, pos_emu, workdir):
    path = workdir / img_entry["file"]
    left, top, width, height = pos_emu or content_area(prs)
    # keep the source aspect ratio inside the target box
    px_w, px_h = img_entry.get("size_px") or (0, 0)
    if px_w and px_h:
        box_ratio, img_ratio = width / height, px_w / px_h
        if img_ratio > box_ratio:
            new_h = int(width / img_ratio)
            top = top + int((height - new_h) / 2)
            height = new_h
        else:
            new_w = int(height * img_ratio)
            left = left + int((width - new_w) / 2)
            width = new_w
    slide.shapes.add_picture(str(path), left, top, width, height)


def place_table(slide, prs, tbl, pos_emu):
    left, top, width, height = pos_emu or content_area(prs)
    shape = slide.shapes.add_table(tbl["n_rows"], tbl["n_cols"], left, top, width, height)
    for ri, row in enumerate(tbl["rows"]):
        for ci, cell_paras in enumerate(row):
            fill_text_frame(shape.table.cell(ri, ci).text_frame, cell_paras)


def strip_existing_slides(prs):
    """Templates usually ship with example slides — the rebuild starts empty."""
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rId = sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)


def rebuild(template_path, manifest_path, mapping_path, out_path):
    prs = Presentation(template_path)
    strip_existing_slides(prs)
    manifest = json.loads(Path(manifest_path).read_text())
    mapping = json.loads(Path(mapping_path).read_text())
    workdir = Path(manifest_path).parent
    drop = set(mapping.get("drop_slides", []))

    # scale factor if slide sizes differ (e.g. 4:3 source onto 16:9 template)
    src_w, src_h = manifest.get("slide_size_emu", [prs.slide_width, prs.slide_height])
    sx, sy = prs.slide_width / src_w, prs.slide_height / src_h

    def scaled(pos):
        if not pos:
            return None
        l, t, w, h = pos
        return Emu(int(l * sx)), Emu(int(t * sy)), Emu(int(w * sx)), Emu(int(h * sy))

    for s in manifest["slides"]:
        n = s["n"]
        if n in drop:
            continue
        new_name = mapping.get("layouts", {}).get(s["layout"]) or mapping.get("default_layout")
        layout = layout_by_name(prs, new_name) if new_name else None
        if layout is None:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            flag(n, f"no mapping for old layout '{s['layout']}' — used fallback '{layout.name}'")
        slide = prs.slides.add_slide(layout)

        bodies = body_placeholders(slide)
        b = 0
        for shape in s["shapes"]:
            ph_type = (shape.get("placeholder") or {}).get("type")
            if ph_type in ("TITLE", "CENTER_TITLE"):
                tp = title_placeholder(slide)
                if tp is not None and shape.get("paragraphs"):
                    fill_text_frame(tp.text_frame, shape["paragraphs"])
                continue
            if "image" in shape:
                # prefer a PICTURE/content placeholder position when available
                place_image(slide, prs, shape["image"], scaled(shape.get("pos_emu")), workdir)
                continue
            if "table" in shape:
                place_table(slide, prs, shape["table"], scaled(shape.get("pos_emu")))
                continue
            if "chart" in shape:
                flag(n, f"chart '{shape.get('chart', {}).get('chart_type')}' needs recreation (data extracted in manifest)")
                continue
            if "group" in shape:
                flag(n, "grouped shape skipped — recreate or paste manually")
                continue
            if shape.get("paragraphs"):
                if b < len(bodies):
                    fill_text_frame(bodies[b].text_frame, shape["paragraphs"])
                    b += 1
                else:
                    from pptx.util import Pt
                    l, t, w, h = scaled(shape.get("pos_emu")) or content_area(prs)
                    box = slide.shapes.add_textbox(l, t, w, h)
                    fill_text_frame(box.text_frame, shape["paragraphs"])
                    box.text_frame.word_wrap = True
                    flag(n, f"text '{shape['paragraphs'][0].get('text','')[:40]}' had no free placeholder — loose textbox")
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    prs.save(out_path)
    print(f"wrote {out_path} ({len(manifest['slides']) - len(drop)} slides)")
    if REVIEW:
        print("\nREVIEW FLAGS:")
        for r in REVIEW:
            print(f"  - {r}")
        Path(out_path).with_suffix(".review.txt").write_text("\n".join(REVIEW))


if __name__ == "__main__":
    if len(sys.argv) != 5:
        sys.exit(__doc__)
    rebuild(*sys.argv[1:])
